"""
PowerPoint生成工具，用于创建和操作PPTX文档。

此模块提供了一组工具函数，用于创建PowerPoint演示文稿、添加幻灯片、插入内容元素、
应用样式和保存文件。基于python-pptx库实现，支持丰富的PPT功能。
"""

import os
import json
from typing import Dict, List, Any, Optional, Tuple, Union
from pathlib import Path

import pptx
from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from strands import tool


@tool
def create_presentation(
    template_path: str = None,
    title: str = None,
    author: str = None,
    output_path: str = None
) -> str:
    """
    创建新的PowerPoint演示文稿，可选择使用模板。

    Args:
        template_path (str, optional): PowerPoint模板文件路径，如果不提供则创建空白演示文稿
        title (str, optional): 演示文稿标题
        author (str, optional): 演示文稿作者
        output_path (str, optional): 输出文件路径，如果不提供则不保存文件
        
    Returns:
        str: JSON格式的结果，包含演示文稿信息和临时文件路径
    """
    try:
        # 创建演示文稿
        if template_path and os.path.exists(template_path):
            prs = Presentation(template_path)
        else:
            prs = Presentation()
        
        # 设置文档属性
        if title or author:
            core_properties = prs.core_properties
            if title:
                core_properties.title = title
            if author:
                core_properties.author = author
        
        # 获取演示文稿信息
        slide_layouts = []
        for idx, layout in enumerate(prs.slide_layouts):
            layout_info = {
                "index": idx,
                "name": layout.name if hasattr(layout, "name") else f"Layout {idx}",
                "placeholder_count": len(layout.placeholders)
            }
            slide_layouts.append(layout_info)
        
        # 保存文件
        temp_path = None
        if output_path:
            # 确保输出目录存在
            output_dir = os.path.dirname(output_path)
            if output_dir and not os.path.exists(output_dir):
                os.makedirs(output_dir)
            
            prs.save(output_path)
            temp_path = output_path
        else:
            # 创建临时文件
            import tempfile
            fd, temp_path = tempfile.mkstemp(suffix='.pptx')
            os.close(fd)
            prs.save(temp_path)
        
        # 构建响应
        response = {
            "status": "success",
            "presentation_info": {
                "slide_count": len(prs.slides),
                "slide_layouts": slide_layouts,
                "title": title,
                "author": author
            },
            "file_path": temp_path
        }
        
        return json.dumps(response, ensure_ascii=False, indent=2)
    
    except Exception as e:
        error_response = {
            "status": "failed",
            "error": str(e),
            "message": "创建演示文稿失败"
        }
        return json.dumps(error_response, ensure_ascii=False)


@tool
def add_slide(
    presentation_path: str,
    layout_index: int = 0,
    title: str = None,
    subtitle: str = None,
    notes: str = None,
    save_path: str = None
) -> str:
    """
    向演示文稿添加新的幻灯片。

    Args:
        presentation_path (str): PowerPoint文件路径
        layout_index (int): 幻灯片布局索引，默认为0（标题幻灯片）
        title (str, optional): 幻灯片标题
        subtitle (str, optional): 幻灯片副标题
        notes (str, optional): 幻灯片备注
        save_path (str, optional): 保存路径，如果不提供则覆盖原文件
        
    Returns:
        str: JSON格式的结果，包含新幻灯片信息
    """
    try:
        # 加载演示文稿
        prs = Presentation(presentation_path)
        
        # 检查布局索引是否有效
        if layout_index < 0 or layout_index >= len(prs.slide_layouts):
            raise ValueError(f"无效的布局索引: {layout_index}，有效范围: 0-{len(prs.slide_layouts)-1}")
        
        # 添加幻灯片
        slide_layout = prs.slide_layouts[layout_index]
        slide = prs.slides.add_slide(slide_layout)
        
        # 设置标题和副标题
        if hasattr(slide, 'shapes') and hasattr(slide.shapes, 'title') and slide.shapes.title:
            if title:
                slide.shapes.title.text = title
        
        # 查找副标题占位符
        if subtitle:
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 1:  # 副标题占位符类型
                    shape.text = subtitle
                    break
        
        # 添加备注
        if notes:
            if not slide.notes_slide:
                slide.notes_slide
            if slide.notes_slide and hasattr(slide.notes_slide, 'notes_text_frame'):
                slide.notes_slide.notes_text_frame.text = notes
        
        # 保存文件
        if save_path:
            # 确保输出目录存在
            output_dir = os.path.dirname(save_path)
            if output_dir and not os.path.exists(output_dir):
                os.makedirs(output_dir)
            prs.save(save_path)
        else:
            prs.save(presentation_path)
        
        # 构建响应
        response = {
            "status": "success",
            "slide_info": {
                "index": len(prs.slides) - 1,
                "layout_index": layout_index,
                "layout_name": slide_layout.name if hasattr(slide_layout, "name") else f"Layout {layout_index}",
                "title": title,
                "subtitle": subtitle,
                "has_notes": notes is not None
            },
            "presentation_info": {
                "slide_count": len(prs.slides),
                "file_path": save_path if save_path else presentation_path
            }
        }
        
        return json.dumps(response, ensure_ascii=False, indent=2)
    
    except Exception as e:
        error_response = {
            "status": "failed",
            "error": str(e),
            "message": "添加幻灯片失败"
        }
        return json.dumps(error_response, ensure_ascii=False)


@tool
def add_text(
    presentation_path: str,
    slide_index: int,
    text: str,
    position: Dict[str, float] = None,
    size: Dict[str, float] = None,
    formatting: Dict[str, Any] = None,
    placeholder_index: int = None,
    save_path: str = None
) -> str:
    """
    向幻灯片添加文本。

    Args:
        presentation_path (str): PowerPoint文件路径
        slide_index (int): 幻灯片索引
        text (str): 要添加的文本
        position (Dict[str, float], optional): 文本框位置，格式: {"left": 值, "top": 值}，单位为英寸
        size (Dict[str, float], optional): 文本框大小，格式: {"width": 值, "height": 值}，单位为英寸
        formatting (Dict[str, Any], optional): 文本格式，可包含:
            - "font_name": 字体名称
            - "font_size": 字体大小（磅）
            - "bold": 是否粗体
            - "italic": 是否斜体
            - "underline": 是否下划线
            - "color": 颜色，格式: [R, G, B]
            - "alignment": 对齐方式，可选值: "left", "center", "right", "justify"
        placeholder_index (int, optional): 占位符索引，如果提供则将文本添加到指定占位符
        save_path (str, optional): 保存路径，如果不提供则覆盖原文件
        
    Returns:
        str: JSON格式的结果
    """
    try:
        # 加载演示文稿
        prs = Presentation(presentation_path)
        
        # 检查幻灯片索引是否有效
        if slide_index < 0 or slide_index >= len(prs.slides):
            raise ValueError(f"无效的幻灯片索引: {slide_index}，有效范围: 0-{len(prs.slides)-1}")
        
        slide = prs.slides[slide_index]
        
        # 处理文本
        if placeholder_index is not None:
            # 添加到占位符
            try:
                placeholder = slide.placeholders[placeholder_index]
                if hasattr(placeholder, 'text'):
                    placeholder.text = text
                else:
                    raise ValueError(f"占位符 {placeholder_index} 不支持文本")
                
                # 应用格式
                if formatting and hasattr(placeholder, 'text_frame'):
                    _apply_text_formatting(placeholder.text_frame, formatting)
                
                shape = placeholder
            except (IndexError, KeyError):
                raise ValueError(f"幻灯片中不存在索引为 {placeholder_index} 的占位符")
        else:
            # 创建文本框
            left = Inches(position['left']) if position and 'left' in position else Inches(1)
            top = Inches(position['top']) if position and 'top' in position else Inches(2)
            width = Inches(size['width']) if size and 'width' in size else Inches(8)
            height = Inches(size['height']) if size and 'height' in size else Inches(1)
            
            shape = slide.shapes.add_textbox(left, top, width, height)
            shape.text_frame.text = text
            
            # 应用格式
            if formatting:
                _apply_text_formatting(shape.text_frame, formatting)
        
        # 保存文件
        if save_path:
            # 确保输出目录存在
            output_dir = os.path.dirname(save_path)
            if output_dir and not os.path.exists(output_dir):
                os.makedirs(output_dir)
            prs.save(save_path)
        else:
            prs.save(presentation_path)
        
        # 构建响应
        response = {
            "status": "success",
            "text_info": {
                "slide_index": slide_index,
                "text": text[:100] + ("..." if len(text) > 100 else ""),
                "shape_id": shape.shape_id if hasattr(shape, 'shape_id') else None,
                "is_placeholder": placeholder_index is not None
            },
            "file_path": save_path if save_path else presentation_path
        }
        
        return json.dumps(response, ensure_ascii=False, indent=2)
    
    except Exception as e:
        error_response = {
            "status": "failed",
            "error": str(e),
            "message": "添加文本失败"
        }
        return json.dumps(error_response, ensure_ascii=False)


@tool
def add_image(
    presentation_path: str,
    slide_index: int,
    image_path: str,
    position: Dict[str, float] = None,
    size: Dict[str, float] = None,
    placeholder_index: int = None,
    save_path: str = None
) -> str:
    """
    向幻灯片添加图片。

    Args:
        presentation_path (str): PowerPoint文件路径
        slide_index (int): 幻灯片索引
        image_path (str): 图片文件路径
        position (Dict[str, float], optional): 图片位置，格式: {"left": 值, "top": 值}，单位为英寸
        size (Dict[str, float], optional): 图片大小，格式: {"width": 值, "height": 值}，单位为英寸
        placeholder_index (int, optional): 占位符索引，如果提供则将图片添加到指定占位符
        save_path (str, optional): 保存路径，如果不提供则覆盖原文件
        
    Returns:
        str: JSON格式的结果
    """
    try:
        # 检查图片文件是否存在
        if not os.path.exists(image_path):
            raise FileNotFoundError(f"图片文件不存在: {image_path}")
        
        # 加载演示文稿
        prs = Presentation(presentation_path)
        
        # 检查幻灯片索引是否有效
        if slide_index < 0 or slide_index >= len(prs.slides):
            raise ValueError(f"无效的幻灯片索引: {slide_index}，有效范围: 0-{len(prs.slides)-1}")
        
        slide = prs.slides[slide_index]
        
        # 处理图片
        if placeholder_index is not None:
            # 添加到占位符
            try:
                placeholder = slide.placeholders[placeholder_index]
                placeholder.insert_picture(image_path)
                shape = placeholder
            except (IndexError, KeyError):
                raise ValueError(f"幻灯片中不存在索引为 {placeholder_index} 的占位符")
            except TypeError:
                raise TypeError(f"占位符 {placeholder_index} 不支持插入图片")
        else:
            # 直接添加图片
            left = Inches(position['left']) if position and 'left' in position else Inches(1)
            top = Inches(position['top']) if position and 'top' in position else Inches(2)
            
            if size and 'width' in size and 'height' in size:
                width = Inches(size['width'])
                height = Inches(size['height'])
                shape = slide.shapes.add_picture(image_path, left, top, width, height)
            else:
                shape = slide.shapes.add_picture(image_path, left, top)
        
        # 保存文件
        if save_path:
            # 确保输出目录存在
            output_dir = os.path.dirname(save_path)
            if output_dir and not os.path.exists(output_dir):
                os.makedirs(output_dir)
            prs.save(save_path)
        else:
            prs.save(presentation_path)
        
        # 构建响应
        response = {
            "status": "success",
            "image_info": {
                "slide_index": slide_index,
                "image_path": image_path,
                "shape_id": shape.shape_id if hasattr(shape, 'shape_id') else None,
                "is_placeholder": placeholder_index is not None
            },
            "file_path": save_path if save_path else presentation_path
        }
        
        return json.dumps(response, ensure_ascii=False, indent=2)
    
    except Exception as e:
        error_response = {
            "status": "failed",
            "error": str(e),
            "message": "添加图片失败"
        }
        return json.dumps(error_response, ensure_ascii=False)


@tool
def add_table(
    presentation_path: str,
    slide_index: int,
    data: List[List[str]],
    position: Dict[str, float] = None,
    size: Dict[str, float] = None,
    formatting: Dict[str, Any] = None,
    has_header: bool = True,
    save_path: str = None
) -> str:
    """
    向幻灯片添加表格。

    Args:
        presentation_path (str): PowerPoint文件路径
        slide_index (int): 幻灯片索引
        data (List[List[str]]): 表格数据，二维数组
        position (Dict[str, float], optional): 表格位置，格式: {"left": 值, "top": 值}，单位为英寸
        size (Dict[str, float], optional): 表格大小，格式: {"width": 值, "height": 值}，单位为英寸
        formatting (Dict[str, Any], optional): 表格格式，可包含:
            - "font_name": 字体名称
            - "font_size": 字体大小（磅）
            - "header_bold": 表头是否粗体
            - "header_bg_color": 表头背景色，格式: [R, G, B]
            - "header_text_color": 表头文字颜色，格式: [R, G, B]
            - "cell_bg_color": 单元格背景色，格式: [R, G, B]
            - "cell_text_color": 单元格文字颜色，格式: [R, G, B]
        has_header (bool): 第一行是否为表头
        save_path (str, optional): 保存路径，如果不提供则覆盖原文件
        
    Returns:
        str: JSON格式的结果
    """
    try:
        # 检查数据有效性
        if not data or not isinstance(data, list) or not all(isinstance(row, list) for row in data):
            raise ValueError("无效的表格数据格式，应为二维数组")
        
        # 加载演示文稿
        prs = Presentation(presentation_path)
        
        # 检查幻灯片索引是否有效
        if slide_index < 0 or slide_index >= len(prs.slides):
            raise ValueError(f"无效的幻灯片索引: {slide_index}，有效范围: 0-{len(prs.slides)-1}")
        
        slide = prs.slides[slide_index]
        
        # 确定表格尺寸
        rows = len(data)
        cols = max(len(row) for row in data) if data else 0
        
        if rows == 0 or cols == 0:
            raise ValueError("表格必须至少包含一行一列")
        
        # 确定表格位置和大小
        left = Inches(position['left']) if position and 'left' in position else Inches(1)
        top = Inches(position['top']) if position and 'top' in position else Inches(2)
        width = Inches(size['width']) if size and 'width' in size else Inches(8)
        height = Inches(size['height']) if size and 'height' in size else Inches(rows * 0.5)
        
        # 创建表格
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # 填充数据
        for row_idx, row_data in enumerate(data):
            for col_idx, cell_text in enumerate(row_data):
                if col_idx < cols:  # 确保不超出列数
                    cell = table.cell(row_idx, col_idx)
                    cell.text = str(cell_text)
        
        # 应用格式
        if formatting:
            _apply_table_formatting(table, formatting, has_header)
        
        # 保存文件
        if save_path:
            # 确保输出目录存在
            output_dir = os.path.dirname(save_path)
            if output_dir and not os.path.exists(output_dir):
                os.makedirs(output_dir)
            prs.save(save_path)
        else:
            prs.save(presentation_path)
        
        # 构建响应
        response = {
            "status": "success",
            "table_info": {
                "slide_index": slide_index,
                "rows": rows,
                "columns": cols,
                "has_header": has_header
            },
            "file_path": save_path if save_path else presentation_path
        }
        
        return json.dumps(response, ensure_ascii=False, indent=2)
    
    except Exception as e:
        error_response = {
            "status": "failed",
            "error": str(e),
            "message": "添加表格失败"
        }
        return json.dumps(error_response, ensure_ascii=False)


@tool
def add_shape(
    presentation_path: str,
    slide_index: int,
    shape_type: str,
    position: Dict[str, float] = None,
    size: Dict[str, float] = None,
    text: str = None,
    formatting: Dict[str, Any] = None,
    save_path: str = None
) -> str:
    """
    向幻灯片添加形状。

    Args:
        presentation_path (str): PowerPoint文件路径
        slide_index (int): 幻灯片索引
        shape_type (str): 形状类型，可选值: "rectangle", "oval", "rounded_rectangle", "diamond", "triangle", "arrow"
        position (Dict[str, float], optional): 形状位置，格式: {"left": 值, "top": 值}，单位为英寸
        size (Dict[str, float], optional): 形状大小，格式: {"width": 值, "height": 值}，单位为英寸
        text (str, optional): 形状中的文本
        formatting (Dict[str, Any], optional): 形状格式，可包含:
            - "fill_color": 填充颜色，格式: [R, G, B]
            - "line_color": 线条颜色，格式: [R, G, B]
            - "line_width": 线条宽度，单位为磅
            - "font_name": 字体名称
            - "font_size": 字体大小（磅）
            - "font_color": 字体颜色，格式: [R, G, B]
            - "text_align": 文本对齐方式，可选值: "left", "center", "right"
        save_path (str, optional): 保存路径，如果不提供则覆盖原文件
        
    Returns:
        str: JSON格式的结果
    """
    try:
        # 加载演示文稿
        prs = Presentation(presentation_path)
        
        # 检查幻灯片索引是否有效
        if slide_index < 0 or slide_index >= len(prs.slides):
            raise ValueError(f"无效的幻灯片索引: {slide_index}，有效范围: 0-{len(prs.slides)-1}")
        
        slide = prs.slides[slide_index]
        
        # 映射形状类型
        shape_types = {
            "rectangle": MSO_SHAPE.RECTANGLE,
            "oval": MSO_SHAPE.OVAL,
            "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
            "diamond": MSO_SHAPE.DIAMOND,
            "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
            "arrow": MSO_SHAPE.RIGHT_ARROW
        }
        
        if shape_type not in shape_types:
            raise ValueError(f"不支持的形状类型: {shape_type}，支持的类型: {', '.join(shape_types.keys())}")
        
        # 确定形状位置和大小
        left = Inches(position['left']) if position and 'left' in position else Inches(2)
        top = Inches(position['top']) if position and 'top' in position else Inches(2)
        width = Inches(size['width']) if size and 'width' in size else Inches(2)
        height = Inches(size['height']) if size and 'height' in size else Inches(1)
        
        # 添加形状
        shape = slide.shapes.add_shape(shape_types[shape_type], left, top, width, height)
        
        # 添加文本
        if text:
            shape.text = text
        
        # 应用格式
        if formatting:
            _apply_shape_formatting(shape, formatting)
        
        # 保存文件
        if save_path:
            # 确保输出目录存在
            output_dir = os.path.dirname(save_path)
            if output_dir and not os.path.exists(output_dir):
                os.makedirs(output_dir)
            prs.save(save_path)
        else:
            prs.save(presentation_path)
        
        # 构建响应
        response = {
            "status": "success",
            "shape_info": {
                "slide_index": slide_index,
                "shape_type": shape_type,
                "shape_id": shape.shape_id if hasattr(shape, 'shape_id') else None,
                "has_text": text is not None
            },
            "file_path": save_path if save_path else presentation_path
        }
        
        return json.dumps(response, ensure_ascii=False, indent=2)
    
    except Exception as e:
        error_response = {
            "status": "failed",
            "error": str(e),
            "message": "添加形状失败"
        }
        return json.dumps(error_response, ensure_ascii=False)


@tool
def get_presentation_info(
    presentation_path: str,
    include_slides: bool = True,
    include_placeholders: bool = False
) -> str:
    """
    获取演示文稿的详细信息。

    Args:
        presentation_path (str): PowerPoint文件路径
        include_slides (bool): 是否包含幻灯片信息
        include_placeholders (bool): 是否包含占位符信息
        
    Returns:
        str: JSON格式的演示文稿信息
    """
    try:
        # 检查文件是否存在
        if not os.path.exists(presentation_path):
            raise FileNotFoundError(f"文件不存在: {presentation_path}")
        
        # 加载演示文稿
        prs = Presentation(presentation_path)
        
        # 基本信息
        info = {
            "file_path": presentation_path,
            "file_name": os.path.basename(presentation_path),
            "file_size": os.path.getsize(presentation_path),
            "slide_count": len(prs.slides),
            "slide_layouts_count": len(prs.slide_layouts)
        }
        
        # 文档属性
        if hasattr(prs, 'core_properties'):
            core_props = prs.core_properties
            info["properties"] = {
                "title": core_props.title if hasattr(core_props, 'title') else None,
                "author": core_props.author if hasattr(core_props, 'author') else None,
                "subject": core_props.subject if hasattr(core_props, 'subject') else None,
                "created": str(core_props.created) if hasattr(core_props, 'created') else None,
                "modified": str(core_props.modified) if hasattr(core_props, 'modified') else None
            }
        
        # 幻灯片布局
        layouts = []
        for idx, layout in enumerate(prs.slide_layouts):
            layout_info = {
                "index": idx,
                "name": layout.name if hasattr(layout, "name") else f"Layout {idx}",
                "placeholder_count": len(layout.placeholders)
            }
            layouts.append(layout_info)
        
        info["slide_layouts"] = layouts
        
        # 幻灯片信息
        if include_slides:
            slides = []
            for idx, slide in enumerate(prs.slides):
                slide_info = {
                    "index": idx,
                    "layout_index": _get_slide_layout_index(slide, prs),
                    "shape_count": len(slide.shapes),
                    "has_notes": hasattr(slide, 'notes_slide') and slide.notes_slide is not None
                }
                
                # 获取标题
                if hasattr(slide, 'shapes') and hasattr(slide.shapes, 'title') and slide.shapes.title:
                    slide_info["title"] = slide.shapes.title.text
                
                # 占位符信息
                if include_placeholders:
                    placeholders = []
                    for p_idx, placeholder in enumerate(slide.placeholders):
                        placeholder_info = {
                            "index": p_idx,
                            "type": placeholder.placeholder_format.type if hasattr(placeholder, 'placeholder_format') else None,
                            "name": placeholder.name if hasattr(placeholder, 'name') else None
                        }
                        placeholders.append(placeholder_info)
                    
                    slide_info["placeholders"] = placeholders
                
                slides.append(slide_info)
            
            info["slides"] = slides
        
        return json.dumps(info, ensure_ascii=False, indent=2)
    
    except Exception as e:
        error_response = {
            "status": "failed",
            "error": str(e),
            "message": "获取演示文稿信息失败"
        }
        return json.dumps(error_response, ensure_ascii=False)


# 辅助函数

def _apply_text_formatting(text_frame, formatting):
    """应用文本格式"""
    if not formatting:
        return
    
    # 应用段落格式
    for paragraph in text_frame.paragraphs:
        # 对齐方式
        if 'alignment' in formatting:
            alignment_map = {
                "left": PP_ALIGN.LEFT,
                "center": PP_ALIGN.CENTER,
                "right": PP_ALIGN.RIGHT,
                "justify": PP_ALIGN.JUSTIFY
            }
            if formatting['alignment'] in alignment_map:
                paragraph.alignment = alignment_map[formatting['alignment']]
        
        # 应用字符格式
        for run in paragraph.runs:
            # 字体名称
            if 'font_name' in formatting:
                run.font.name = formatting['font_name']
            
            # 字体大小
            if 'font_size' in formatting:
                run.font.size = Pt(formatting['font_size'])
            
            # 粗体
            if 'bold' in formatting:
                run.font.bold = formatting['bold']
            
            # 斜体
            if 'italic' in formatting:
                run.font.italic = formatting['italic']
            
            # 下划线
            if 'underline' in formatting:
                run.font.underline = formatting['underline']
            
            # 颜色
            if 'color' in formatting and isinstance(formatting['color'], list) and len(formatting['color']) == 3:
                r, g, b = formatting['color']
                run.font.color.rgb = RGBColor(r, g, b)


def _apply_table_formatting(table, formatting, has_header):
    """应用表格格式"""
    if not formatting:
        return
    
    # 遍历所有单元格
    for row_idx, row in enumerate(table.rows):
        is_header = row_idx == 0 and has_header
        
        for cell in row.cells:
            # 设置单元格背景色
            if is_header and 'header_bg_color' in formatting:
                cell.fill.solid()
                r, g, b = formatting['header_bg_color']
                cell.fill.fore_color.rgb = RGBColor(r, g, b)
            elif not is_header and 'cell_bg_color' in formatting:
                cell.fill.solid()
                r, g, b = formatting['cell_bg_color']
                cell.fill.fore_color.rgb = RGBColor(r, g, b)
            
            # 应用文本格式
            for paragraph in cell.text_frame.paragraphs:
                # 字体名称
                if 'font_name' in formatting:
                    for run in paragraph.runs:
                        run.font.name = formatting['font_name']
                
                # 字体大小
                if 'font_size' in formatting:
                    for run in paragraph.runs:
                        run.font.size = Pt(formatting['font_size'])
                
                # 表头粗体
                if is_header and 'header_bold' in formatting and formatting['header_bold']:
                    for run in paragraph.runs:
                        run.font.bold = True
                
                # 文字颜色
                if is_header and 'header_text_color' in formatting:
                    for run in paragraph.runs:
                        r, g, b = formatting['header_text_color']
                        run.font.color.rgb = RGBColor(r, g, b)
                elif not is_header and 'cell_text_color' in formatting:
                    for run in paragraph.runs:
                        r, g, b = formatting['cell_text_color']
                        run.font.color.rgb = RGBColor(r, g, b)


def _apply_shape_formatting(shape, formatting):
    """应用形状格式"""
    if not formatting:
        return
    
    # 填充颜色
    if 'fill_color' in formatting:
        shape.fill.solid()
        r, g, b = formatting['fill_color']
        shape.fill.fore_color.rgb = RGBColor(r, g, b)
    
    # 线条颜色
    if 'line_color' in formatting:
        r, g, b = formatting['line_color']
        shape.line.color.rgb = RGBColor(r, g, b)
    
    # 线条宽度
    if 'line_width' in formatting:
        shape.line.width = Pt(formatting['line_width'])
    
    # 文本格式
    if hasattr(shape, 'text_frame'):
        text_formatting = {}
        
        if 'font_name' in formatting:
            text_formatting['font_name'] = formatting['font_name']
        
        if 'font_size' in formatting:
            text_formatting['font_size'] = formatting['font_size']
        
        if 'font_color' in formatting:
            text_formatting['color'] = formatting['font_color']
        
        if 'text_align' in formatting:
            text_formatting['alignment'] = formatting['text_align']
        
        _apply_text_formatting(shape.text_frame, text_formatting)


def _get_slide_layout_index(slide, presentation):
    """获取幻灯片布局索引"""
    if not hasattr(slide, 'slide_layout'):
        return -1
    
    for idx, layout in enumerate(presentation.slide_layouts):
        if slide.slide_layout == layout:
            return idx
    
    return -1