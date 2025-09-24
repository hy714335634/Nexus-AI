from typing import Dict, List, Any, Optional, Union, Tuple
import json
import os
import base64
import re
import logging
import hashlib
from datetime import datetime
import tempfile
from pathlib import Path
import boto3
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from io import BytesIO
import xml.etree.ElementTree as ET

from strands import tool

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

# Constants
AWS_ICON_BASE_URL = "https://d1.awsstatic.com/icons/aws-icons/"
AWS_ICON_FALLBACK_URL = "https://d1.awsstatic.com/icons/aws-icons/General/AWS-Cloud.svg"
DEFAULT_TEMPLATE_PATH = os.path.join(os.path.dirname(__file__), "templates", "aws_architecture_template.pptx")

# Cache directory for AWS icons and templates
CACHE_DIR = os.path.join(os.path.expanduser("~"), ".cache", "aws_architect")
ICONS_CACHE_DIR = os.path.join(CACHE_DIR, "icons")
TEMPLATES_CACHE_DIR = os.path.join(CACHE_DIR, "templates")

# Create cache directories if they don't exist
os.makedirs(ICONS_CACHE_DIR, exist_ok=True)
os.makedirs(TEMPLATES_CACHE_DIR, exist_ok=True)

# AWS service to icon mapping
AWS_SERVICE_ICONS = {
    "ec2": "Compute/Amazon-EC2.svg",
    "s3": "Storage/Amazon-S3.svg",
    "rds": "Database/Amazon-RDS.svg",
    "lambda": "Compute/AWS-Lambda.svg",
    "dynamodb": "Database/Amazon-DynamoDB.svg",
    "vpc": "Networking-Content-Delivery/Amazon-VPC.svg",
    "cloudfront": "Networking-Content-Delivery/Amazon-CloudFront.svg",
    "route53": "Networking-Content-Delivery/Amazon-Route-53.svg",
    "elb": "Networking-Content-Delivery/Elastic-Load-Balancing.svg",
    "alb": "Networking-Content-Delivery/Elastic-Load-Balancing_Application-Load-Balancer.svg",
    "nlb": "Networking-Content-Delivery/Elastic-Load-Balancing_Network-Load-Balancer.svg",
    "api gateway": "Networking-Content-Delivery/Amazon-API-Gateway.svg",
    "cloudwatch": "Management-Governance/Amazon-CloudWatch.svg",
    "iam": "Security-Identity-Compliance/AWS-Identity-and-Access-Management.svg",
    "sns": "Application-Integration/Amazon-SNS.svg",
    "sqs": "Application-Integration/Amazon-SQS.svg",
    "aurora": "Database/Amazon-Aurora.svg",
    "elasticache": "Database/Amazon-ElastiCache.svg",
    "redshift": "Database/Amazon-Redshift.svg",
    "efs": "Storage/Amazon-EFS.svg",
    "waf": "Security-Identity-Compliance/AWS-WAF.svg",
    "shield": "Security-Identity-Compliance/AWS-Shield.svg",
    "cognito": "Security-Identity-Compliance/Amazon-Cognito.svg",
    "kms": "Security-Identity-Compliance/AWS-Key-Management-Service.svg",
    "cloudtrail": "Management-Governance/AWS-CloudTrail.svg",
    "config": "Management-Governance/AWS-Config.svg",
    "cloudformation": "Management-Governance/AWS-CloudFormation.svg",
    "step functions": "Application-Integration/AWS-Step-Functions.svg",
    "eventbridge": "Application-Integration/Amazon-EventBridge.svg",
    "ecs": "Compute/Amazon-ECS.svg",
    "eks": "Compute/Amazon-EKS.svg",
    "fargate": "Compute/AWS-Fargate.svg",
    "elastic beanstalk": "Compute/AWS-Elastic-Beanstalk.svg",
    "internet gateway": "Networking-Content-Delivery/Amazon-VPC_Internet-Gateway.svg",
    "nat gateway": "Networking-Content-Delivery/Amazon-VPC_NAT-Gateway.svg",
    "subnet": "Networking-Content-Delivery/Amazon-VPC_Subnet.svg",
    "security group": "Networking-Content-Delivery/Amazon-VPC_Security-Group.svg",
    "direct connect": "Networking-Content-Delivery/AWS-Direct-Connect.svg",
    "transit gateway": "Networking-Content-Delivery/AWS-Transit-Gateway.svg"
}

# AWS Best Practices
AWS_BEST_PRACTICES = {
    "security": [
        "Use security groups and NACLs to control traffic",
        "Implement least privilege IAM policies",
        "Enable encryption for data at rest and in transit",
        "Use AWS WAF for web application protection",
        "Implement multi-factor authentication for IAM users",
        "Use AWS Shield for DDoS protection",
        "Regularly audit and rotate access keys",
        "Use AWS Config for compliance monitoring"
    ],
    "reliability": [
        "Deploy across multiple Availability Zones",
        "Use Auto Scaling for handling load changes",
        "Implement health checks and self-healing",
        "Use managed services when possible",
        "Design for failure and test recovery procedures",
        "Use Amazon Route 53 for DNS failover",
        "Implement backup and restore strategies",
        "Use AWS Backup for centralized backup management"
    ],
    "performance": [
        "Use CloudFront for content delivery",
        "Implement caching strategies",
        "Use read replicas for database read scaling",
        "Choose the right instance types for workloads",
        "Use ElastiCache for in-memory caching",
        "Implement asynchronous processing for non-critical tasks",
        "Use S3 Transfer Acceleration for fast file transfers",
        "Optimize database queries and indexing"
    ],
    "cost": [
        "Right-size resources based on actual needs",
        "Use Auto Scaling to match capacity with demand",
        "Implement lifecycle policies for S3 objects",
        "Use Spot Instances for non-critical workloads",
        "Use Reserved Instances for predictable workloads",
        "Implement tagging strategy for cost allocation",
        "Use AWS Cost Explorer and Budgets for monitoring",
        "Optimize storage choices based on access patterns"
    ],
    "operational_excellence": [
        "Use Infrastructure as Code (CloudFormation, CDK)",
        "Implement comprehensive monitoring and alerting",
        "Automate routine operational tasks",
        "Use AWS Systems Manager for operational management",
        "Implement proper tagging for resource management",
        "Use AWS CloudTrail for API activity monitoring",
        "Document operational procedures",
        "Regularly test disaster recovery procedures"
    ]
}

# VPC service classification
VPC_INTERNAL_SERVICES = [
    "ec2", "rds", "elasticache", "efs", "lambda", "ecs", "eks", "fargate", 
    "aurora", "redshift", "elastic beanstalk", "emr", "elasticsearch", 
    "opensearch", "msk", "kafka", "documentdb", "neptune", "memorydb"
]

VPC_EXTERNAL_SERVICES = [
    "s3", "dynamodb", "cloudfront", "route53", "sns", "sqs", "eventbridge",
    "step functions", "api gateway", "cognito", "kms", "iam", "cloudwatch",
    "cloudtrail", "config", "waf", "shield", "glacier", "athena", "quicksight"
]

@tool
def generate_ppt_architecture(architecture_components: List[Dict[str, Any]], 
                              template_path: Optional[str] = None, 
                              diagram_title: str = "AWS Architecture Diagram") -> str:
    """Generate PowerPoint presentation diagram based on architecture components.

    Args:
        architecture_components (List[Dict[str, Any]]): List of architecture components with their relationships
            Each component should have:
            - id: Unique identifier
            - type: AWS service type
            - name: Display name
            - connections: List of component IDs this component connects to
            - vpc (optional): VPC identifier if the component is inside a VPC
        template_path (Optional[str]): Path to PowerPoint template file (.pptx)
        diagram_title (str): Title of the diagram
        
    Returns:
        str: JSON string containing information about the generated PowerPoint file
    """
    if not architecture_components:
        return json.dumps({
            "status": "error",
            "error": "Empty architecture components",
            "message": "Please provide a list of architecture components."
        })
    
    try:
        # Use the provided template or default template
        ppt_template = template_path if template_path and os.path.exists(template_path) else DEFAULT_TEMPLATE_PATH
        
        # Create a temporary file for the output
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_file:
            output_path = temp_file.name
        
        # Create presentation from template or create new if template doesn't exist
        try:
            prs = Presentation(ppt_template)
            logger.info(f"Using template: {ppt_template}")
        except Exception as e:
            logger.warning(f"Failed to use template {ppt_template}: {e}. Creating new presentation.")
            prs = Presentation()
        
        # Create title slide
        title_slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = diagram_title
        subtitle.text = f"Generated on {datetime.now().strftime('%Y-%m-%d')}"
        
        # Create architecture diagram slide
        diagram_slide_layout = prs.slide_layouts[5]  # Blank slide layout
        diagram_slide = prs.slides.add_slide(diagram_slide_layout)
        
        # Add slide title
        title_shape = diagram_slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.5))
        title_frame = title_shape.text_frame
        title_paragraph = title_frame.add_paragraph()
        title_paragraph.text = "Architecture Overview"
        title_paragraph.font.size = Pt(24)
        title_paragraph.font.bold = True
        
        # Track VPCs for grouping
        vpcs = [c for c in architecture_components if c.get("type", "").lower() == "vpc"]
        vpc_components = {}
        
        for vpc in vpcs:
            vpc_id = vpc.get("id")
            vpc_components[vpc_id] = [c for c in architecture_components if c.get("vpc") == vpc_id]
        
        # Calculate layout
        # This is a simplified layout - in a real implementation, you would use a more sophisticated algorithm
        components_by_id = {c.get("id"): c for c in architecture_components if c.get("id")}
        
        # Add VPC containers
        vpc_shapes = {}
        vpc_y_position = 1.5  # Starting Y position in inches
        
        for i, vpc in enumerate(vpcs):
            vpc_id = vpc.get("id")
            vpc_name = vpc.get("name", f"VPC {i+1}")
            
            # Calculate VPC size based on components inside it
            vpc_component_count = len(vpc_components.get(vpc_id, []))
            vpc_width = max(3, min(9, vpc_component_count * 1.2))
            vpc_height = max(2, min(5, vpc_component_count * 0.8))
            
            # Position VPC - alternate between left and right side
            vpc_x = 1 if i % 2 == 0 else 10 - vpc_width
            
            # Create VPC shape
            vpc_shape = diagram_slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 
                Inches(vpc_x), 
                Inches(vpc_y_position), 
                Inches(vpc_width), 
                Inches(vpc_height)
            )
            
            # Style VPC shape
            vpc_shape.fill.solid()
            vpc_shape.fill.fore_color.rgb = RGBColor(240, 248, 255)  # Light blue
            vpc_shape.line.color.rgb = RGBColor(30, 136, 229)  # Medium blue
            vpc_shape.line.width = Pt(2)
            
            # Add VPC name
            text_frame = vpc_shape.text_frame
            text_frame.text = vpc_name
            text_frame.paragraphs[0].font.bold = True
            
            # Store VPC shape for later reference
            vpc_shapes[vpc_id] = {
                "shape": vpc_shape,
                "x": vpc_x,
                "y": vpc_y_position,
                "width": vpc_width,
                "height": vpc_height
            }
            
            # Update Y position for next VPC
            vpc_y_position += vpc_height + 0.5
        
        # Add components to VPCs
        component_shapes = {}
        
        for vpc_id, components in vpc_components.items():
            if vpc_id not in vpc_shapes:
                continue
                
            vpc_info = vpc_shapes[vpc_id]
            vpc_x = vpc_info["x"]
            vpc_y = vpc_info["y"]
            vpc_width = vpc_info["width"]
            vpc_height = vpc_info["height"]
            
            # Calculate grid layout inside VPC
            component_count = len(components)
            cols = min(3, component_count)
            rows = (component_count + cols - 1) // cols
            
            cell_width = vpc_width / cols
            cell_height = vpc_height / (rows + 0.5)  # Add some padding for VPC title
            
            for i, component in enumerate(components):
                component_id = component.get("id")
                component_name = component.get("name", component_id)
                component_type = component.get("type", "Generic").lower()
                
                row = i // cols
                col = i % cols
                
                # Calculate position within VPC
                x = vpc_x + col * cell_width + cell_width * 0.1
                y = vpc_y + row * cell_height + cell_height * 0.5  # Add offset for VPC title
                
                # Add component shape
                component_shape = add_aws_service_icon(
                    diagram_slide, 
                    component_type, 
                    component_name, 
                    x, 
                    y, 
                    cell_width * 0.8, 
                    cell_height * 0.6
                )
                
                # Store component shape for connections
                component_shapes[component_id] = {
                    "shape": component_shape,
                    "x": x + (cell_width * 0.8 / 2),
                    "y": y + (cell_height * 0.6 / 2)
                }
        
        # Add components not in VPCs
        non_vpc_components = [c for c in architecture_components if not c.get("vpc") and c.get("type", "").lower() != "vpc"]
        
        if non_vpc_components:
            # Calculate grid layout for non-VPC components
            component_count = len(non_vpc_components)
            cols = min(4, component_count)
            rows = (component_count + cols - 1) // cols
            
            start_y = vpc_y_position  # Start below the last VPC
            cell_width = 9 / cols
            cell_height = 1.2
            
            for i, component in enumerate(non_vpc_components):
                component_id = component.get("id")
                component_name = component.get("name", component_id)
                component_type = component.get("type", "Generic").lower()
                
                row = i // cols
                col = i % cols
                
                # Calculate position
                x = 1 + col * cell_width
                y = start_y + row * cell_height
                
                # Add component shape
                component_shape = add_aws_service_icon(
                    diagram_slide, 
                    component_type, 
                    component_name, 
                    x, 
                    y, 
                    cell_width * 0.8, 
                    cell_height * 0.6
                )
                
                # Store component shape for connections
                component_shapes[component_id] = {
                    "shape": component_shape,
                    "x": x + (cell_width * 0.8 / 2),
                    "y": y + (cell_height * 0.6 / 2)
                }
        
        # Add connections between components
        for component in architecture_components:
            source_id = component.get("id")
            if source_id not in component_shapes:
                continue
                
            source_info = component_shapes[source_id]
            connections = component.get("connections", [])
            
            for target_id in connections:
                if target_id in component_shapes:
                    target_info = component_shapes[target_id]
                    
                    # Add connector line
                    connector = diagram_slide.shapes.add_connector(
                        MSO_SHAPE.LINE_STRAIGHT,
                        Inches(source_info["x"]),
                        Inches(source_info["y"]),
                        Inches(target_info["x"]),
                        Inches(target_info["y"])
                    )
                    
                    connector.line.color.rgb = RGBColor(128, 128, 128)
                    connector.line.width = Pt(1.5)
        
        # Add legend
        legend_y = vpc_y_position + (0 if not non_vpc_components else (rows * cell_height + 0.5))
        add_legend(diagram_slide, legend_y)
        
        # Create component details slide
        details_slide_layout = prs.slide_layouts[5]  # Blank slide layout
        details_slide = prs.slides.add_slide(details_slide_layout)
        
        # Add slide title
        title_shape = details_slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.5))
        title_frame = title_shape.text_frame
        title_paragraph = title_frame.add_paragraph()
        title_paragraph.text = "Component Details"
        title_paragraph.font.size = Pt(24)
        title_paragraph.font.bold = True
        
        # Group components by type
        components_by_type = {}
        for component in architecture_components:
            component_type = component.get("type", "Other")
            if component_type not in components_by_type:
                components_by_type[component_type] = []
            components_by_type[component_type].append(component)
        
        # Add component details table
        table_y = 1.2
        for component_type, components in components_by_type.items():
            # Skip if too many component types to fit on one slide
            if table_y > 6:
                # Create additional details slide
                details_slide = prs.slides.add_slide(details_slide_layout)
                title_shape = details_slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.5))
                title_frame = title_shape.text_frame
                title_paragraph = title_frame.add_paragraph()
                title_paragraph.text = "Component Details (Continued)"
                title_paragraph.font.size = Pt(24)
                title_paragraph.font.bold = True
                table_y = 1.2
            
            # Add component type heading
            type_shape = details_slide.shapes.add_textbox(Inches(0.5), Inches(table_y), Inches(9), Inches(0.3))
            type_frame = type_shape.text_frame
            type_paragraph = type_frame.add_paragraph()
            type_paragraph.text = component_type
            type_paragraph.font.size = Pt(18)
            type_paragraph.font.bold = True
            
            table_y += 0.4
            
            # Add table with component details
            rows = len(components) + 1  # +1 for header row
            cols = 3
            table_width = 9
            table_height = rows * 0.3
            
            table = details_slide.shapes.add_table(
                rows, cols,
                Inches(0.5), Inches(table_y),
                Inches(table_width), Inches(table_height)
            ).table
            
            # Set column widths
            table.columns[0].width = Inches(2)
            table.columns[1].width = Inches(4)
            table.columns[2].width = Inches(3)
            
            # Add header row
            table.cell(0, 0).text = "Name"
            table.cell(0, 1).text = "Description"
            table.cell(0, 2).text = "Connections"
            
            # Style header row
            for col in range(cols):
                cell = table.cell(0, col)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(230, 230, 230)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.bold = True
            
            # Add component rows
            for i, component in enumerate(components):
                row = i + 1
                component_name = component.get("name", component.get("id", "Unknown"))
                component_description = component.get("description", "No description provided")
                
                # Get connections
                connection_names = []
                for conn_id in component.get("connections", []):
                    conn_component = next((c for c in architecture_components if c.get("id") == conn_id), None)
                    if conn_component:
                        connection_names.append(conn_component.get("name", conn_id))
                
                connections_text = ", ".join(connection_names) if connection_names else "None"
                
                table.cell(row, 0).text = component_name
                table.cell(row, 1).text = component_description
                table.cell(row, 2).text = connections_text
            
            table_y += table_height + 0.3
        
        # Save the presentation
        prs.save(output_path)
        
        result = {
            "status": "success",
            "message": "PowerPoint presentation generated successfully",
            "file_path": output_path,
            "component_count": len(architecture_components),
            "vpc_count": len(vpcs),
            "slide_count": len(prs.slides)
        }
        
        return json.dumps(result)
    except Exception as e:
        logger.error(f"Error generating PowerPoint diagram: {str(e)}")
        return json.dumps({
            "status": "error",
            "error": "Failed to generate PowerPoint diagram",
            "message": str(e)
        })

def add_aws_service_icon(slide, service_type: str, service_name: str, x: float, y: float, width: float, height: float):
    """Add an AWS service icon to the slide."""
    # Normalize service type
    service_type = service_type.lower().strip()
    
    # Create shape for the icon
    icon_shape = slide.shapes.add_picture(
        get_aws_icon_path(service_type),
        Inches(x),
        Inches(y),
        Inches(width * 0.8)
    )
    
    # Add text label below the icon
    text_y = y + height * 0.7
    text_shape = slide.shapes.add_textbox(
        Inches(x - width * 0.1),
        Inches(text_y),
        Inches(width),
        Inches(height * 0.3)
    )
    
    text_frame = text_shape.text_frame
    text_frame.text = service_name
    text_frame.paragraphs[0].alignment = 1  # Center alignment
    text_frame.paragraphs[0].font.size = Pt(10)
    
    # Return the icon shape for connection references
    return icon_shape

def get_aws_icon_path(service_type: str) -> str:
    """Get the path to the AWS service icon, downloading it if necessary."""
    # Normalize service type
    service_type = service_type.lower().strip()
    
    # Check if we have a mapping for this service
    icon_filename = AWS_SERVICE_ICONS.get(service_type)
    
    if not icon_filename:
        # Try to find a partial match
        for key, value in AWS_SERVICE_ICONS.items():
            if key in service_type or service_type in key:
                icon_filename = value
                break
    
    # If still no match, use a generic AWS icon
    if not icon_filename:
        icon_filename = "General/AWS-Cloud.svg"
    
    # Create cache path
    cache_path = os.path.join(ICONS_CACHE_DIR, icon_filename.replace("/", "_"))
    
    # If icon doesn't exist in cache, download it
    if not os.path.exists(cache_path):
        try:
            # In a real implementation, this would download the icon from AWS
            # For now, we'll create a placeholder icon
            create_placeholder_icon(cache_path, service_type)
        except Exception as e:
            logger.error(f"Failed to create icon for {service_type}: {e}")
            # Create a simple placeholder
            create_simple_placeholder(cache_path, service_type)
    
    return cache_path

def create_placeholder_icon(path: str, service_type: str):
    """Create a placeholder icon for the service."""
    # Create a simple SVG icon
    svg_content = f"""
    <svg width="64" height="64" xmlns="http://www.w3.org/2000/svg">
        <rect width="64" height="64" fill="#FF9900" rx="5" ry="5"/>
        <text x="32" y="32" font-family="Arial" font-size="10" text-anchor="middle" fill="white" dominant-baseline="middle">
            {service_type.upper()[:10]}
        </text>
    </svg>
    """
    
    # Ensure directory exists
    os.makedirs(os.path.dirname(path), exist_ok=True)
    
    # Write SVG content to file
    with open(path, 'w') as f:
        f.write(svg_content)

def create_simple_placeholder(path: str, service_type: str):
    """Create a very simple placeholder when SVG creation fails."""
    # Create a directory for the file if it doesn't exist
    os.makedirs(os.path.dirname(path), exist_ok=True)
    
    # Create a simple text file that PowerPoint can still use as a placeholder
    with open(path, 'w') as f:
        f.write(f"AWS {service_type}")

def add_legend(slide, y_position: float):
    """Add a legend to the slide."""
    legend_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5),
        Inches(y_position),
        Inches(9),
        Inches(0.8)
    )
    
    # Style legend shape
    legend_shape.fill.solid()
    legend_shape.fill.fore_color.rgb = RGBColor(245, 245, 245)
    legend_shape.line.color.rgb = RGBColor(200, 200, 200)
    
    # Add legend title
    legend_title = slide.shapes.add_textbox(
        Inches(0.7),
        Inches(y_position + 0.1),
        Inches(2),
        Inches(0.3)
    )
    legend_title.text_frame.text = "Legend:"
    legend_title.text_frame.paragraphs[0].font.bold = True
    
    # Add legend items
    legend_items = [
        ("VPC", "Virtual Private Cloud"),
        ("→", "Connection between services"),
        ("AWS Icons", "Represent AWS services")
    ]
    
    for i, (item, description) in enumerate(legend_items):
        item_x = 2 + i * 2.5
        
        # Add item name
        item_shape = slide.shapes.add_textbox(
            Inches(item_x),
            Inches(y_position + 0.1),
            Inches(1),
            Inches(0.3)
        )
        item_shape.text_frame.text = item
        item_shape.text_frame.paragraphs[0].font.bold = True
        
        # Add item description
        desc_shape = slide.shapes.add_textbox(
            Inches(item_x + 0.8),
            Inches(y_position + 0.1),
            Inches(1.5),
            Inches(0.3)
        )
        desc_shape.text_frame.text = description

@tool
def validate_aws_architecture(architecture_components: List[Dict[str, Any]]) -> str:
    """Validate AWS architecture design for best practices and common issues.

    Args:
        architecture_components (List[Dict[str, Any]]): List of architecture components with their relationships
            Each component should have:
            - id: Unique identifier
            - type: AWS service type
            - name: Display name
            - connections: List of component IDs this component connects to
            - vpc (optional): VPC identifier if the component is inside a VPC
        
    Returns:
        str: JSON string containing validation results, recommendations and issues found
    """
    if not architecture_components:
        return json.dumps({
            "status": "error",
            "error": "Empty architecture components",
            "message": "Please provide a list of architecture components."
        })
    
    try:
        # Initialize validation results
        validation_results = {
            "status": "success",
            "summary": {
                "component_count": len(architecture_components),
                "critical_issues": 0,
                "warning_issues": 0,
                "info_issues": 0,
                "best_practices_met": 0
            },
            "issues": [],
            "best_practices_met": [],
            "recommendations": []
        }
        
        # Helper function to add an issue
        def add_issue(severity: str, title: str, description: str, affected_components: List[str] = None, recommendation: str = None):
            issue = {
                "severity": severity,
                "title": title,
                "description": description
            }
            
            if affected_components:
                issue["affected_components"] = affected_components
                
            if recommendation:
                issue["recommendation"] = recommendation
                
            validation_results["issues"].append(issue)
            
            if severity == "critical":
                validation_results["summary"]["critical_issues"] += 1
            elif severity == "warning":
                validation_results["summary"]["warning_issues"] += 1
            elif severity == "info":
                validation_results["summary"]["info_issues"] += 1
        
        # Helper function to add a best practice that's being followed
        def add_best_practice(category: str, title: str, description: str):
            validation_results["best_practices_met"].append({
                "category": category,
                "title": title,
                "description": description
            })
            validation_results["summary"]["best_practices_met"] += 1
        
        # Helper function to add a recommendation
        def add_recommendation(category: str, title: str, description: str, importance: str = "medium"):
            validation_results["recommendations"].append({
                "category": category,
                "title": title,
                "description": description,
                "importance": importance
            })
        
        # Track components by ID for easier reference
        components_by_id = {c.get("id"): c for c in architecture_components if c.get("id")}
        
        # 1. Check for VPC usage
        vpcs = [c for c in architecture_components if c.get("type", "").lower() == "vpc"]
        
        if not vpcs:
            add_issue(
                "warning",
                "No VPC defined",
                "Architecture does not include a Virtual Private Cloud (VPC), which is recommended for network isolation and security.",
                recommendation="Add a VPC to isolate your resources and control network traffic."
            )
        else:
            add_best_practice(
                "security",
                "VPC Usage",
                "Architecture includes VPC for network isolation and security."
            )
        
        # 2. Check for proper VPC service placement
        misplaced_services = []
        
        for component in architecture_components:
            component_type = component.get("type", "").lower()
            component_name = component.get("name", component.get("id", "Unknown"))
            vpc_id = component.get("vpc")
            
            # Check if service that should be in VPC is not in a VPC
            if component_type in VPC_INTERNAL_SERVICES and not vpc_id:
                misplaced_services.append({
                    "name": component_name,
                    "type": component_type,
                    "issue": "Should be inside a VPC"
                })
            
            # Check if service that should be outside VPC is in a VPC
            if component_type in VPC_EXTERNAL_SERVICES and vpc_id:
                misplaced_services.append({
                    "name": component_name,
                    "type": component_type,
                    "issue": "Should be outside a VPC"
                })
        
        if misplaced_services:
            affected_components = [f"{s['name']} ({s['issue']})" for s in misplaced_services]
            add_issue(
                "warning",
                "Incorrect VPC service placement",
                "Some services are not correctly placed with respect to VPC boundaries.",
                affected_components=affected_components,
                recommendation="Review the placement of services in relation to VPC boundaries according to AWS best practices."
            )
        else:
            add_best_practice(
                "architecture",
                "Proper VPC Service Placement",
                "All services are correctly placed with respect to VPC boundaries."
            )
        
        # 3. Check for security components
        security_components = [c for c in architecture_components if c.get("type", "").lower() in [
            "security group", "nacl", "waf", "shield", "iam", "kms", "secrets manager", 
            "certificate manager", "cognito", "guardduty", "inspector", "macie"
        ]]
        
        if not security_components:
            add_issue(
                "warning",
                "Missing security components",
                "Architecture does not include explicit security components like Security Groups, NACLs, WAF, etc.",
                recommendation="Add appropriate security components to protect your application and data."
            )
        else:
            add_best_practice(
                "security",
                "Security Components",
                "Architecture includes dedicated security components for protection."
            )
        
        # 4. Check for high availability
        has_multi_az = False
        has_load_balancer = False
        has_auto_scaling = False
        
        for component in architecture_components:
            component_description = component.get("description", "").lower()
            component_type = component.get("type", "").lower()
            
            if "multi-az" in component_description or "multiple az" in component_description:
                has_multi_az = True
                
            if component_type in ["elb", "alb", "nlb", "load balancer"]:
                has_load_balancer = True
                
            if "auto scaling" in component_type or "autoscaling" in component_type:
                has_auto_scaling = True
        
        if not has_multi_az and not has_load_balancer:
            add_issue(
                "warning",
                "Limited high availability",
                "Architecture does not appear to implement high availability across multiple AZs or use load balancers.",
                recommendation="Consider deploying resources across multiple AZs and using load balancers for high availability."
            )
        else:
            add_best_practice(
                "reliability",
                "High Availability Design",
                "Architecture implements high availability principles."
            )
        
        # 5. Check for database components
        database_components = [c for c in architecture_components if c.get("type", "").lower() in [
            "rds", "aurora", "dynamodb", "documentdb", "neptune", "redshift", "timestream",
            "keyspaces", "qldb", "elasticache", "memorydb", "database"
        ]]
        
        if database_components:
            # Check for database backup/redundancy mentions
            has_backup_mention = False
            for component in database_components:
                description = component.get("description", "").lower()
                if any(term in description for term in ["backup", "redundancy", "multi-az", "replica"]):
                    has_backup_mention = True
                    break
            
            if not has_backup_mention:
                affected_components = [c.get("name", c.get("id", "Unknown")) for c in database_components]
                add_issue(
                    "info",
                    "Database backup/redundancy not specified",
                    "Architecture includes databases but doesn't explicitly mention backup or redundancy strategies.",
                    affected_components=affected_components,
                    recommendation="Implement automated backups and consider Multi-AZ deployments or read replicas for databases."
                )
        
        # 6. Check for monitoring and logging
        monitoring_components = [c for c in architecture_components if c.get("type", "").lower() in [
            "cloudwatch", "cloudtrail", "config", "x-ray", "managed grafana", "managed prometheus",
            "monitoring", "logs", "logging"
        ]]
        
        if not monitoring_components:
            add_issue(
                "info",
                "Missing monitoring components",
                "Architecture does not include explicit monitoring or logging components.",
                recommendation="Add CloudWatch for monitoring, CloudTrail for API activity logging, and consider other observability services."
            )
        else:
            add_best_practice(
                "operational_excellence",
                "Monitoring and Logging",
                "Architecture includes components for monitoring and logging."
            )
        
        # 7. Check for content delivery optimization
        has_s3 = any(c.get("type", "").lower() == "s3" for c in architecture_components)
        has_cloudfront = any(c.get("type", "").lower() == "cloudfront" for c in architecture_components)
        
        if has_s3 and not has_cloudfront:
            add_recommendation(
                "performance",
                "Add CloudFront for content delivery",
                "Consider using CloudFront CDN with your S3 buckets for improved performance and reduced latency.",
                "medium"
            )
        
        # 8. Check for proper networking components in VPC
        for vpc in vpcs:
            vpc_id = vpc.get("id")
            vpc_name = vpc.get("name", vpc_id)
            
            # Check for internet gateway if there are public-facing components
            has_public_subnet = False
            has_internet_gateway = False
            has_nat_gateway = False
            
            for component in architecture_components:
                if component.get("vpc") == vpc_id:
                    component_type = component.get("type", "").lower()
                    component_description = component.get("description", "").lower()
                    
                    if "public subnet" in component_description or "public-facing" in component_description:
                        has_public_subnet = True
                    
                    if component_type in ["internet gateway", "igw"]:
                        has_internet_gateway = True
                    
                    if component_type in ["nat gateway", "nat"]:
                        has_nat_gateway = True
            
            if has_public_subnet and not has_internet_gateway:
                add_issue(
                    "warning",
                    f"Missing Internet Gateway in {vpc_name}",
                    f"VPC {vpc_name} has public-facing components but no Internet Gateway.",
                    recommendation="Add an Internet Gateway to allow communication between instances in your VPC and the internet."
                )
            
            if has_public_subnet and not has_nat_gateway:
                add_issue(
                    "info",
                    f"Consider adding NAT Gateway in {vpc_name}",
                    f"VPC {vpc_name} has public-facing components but no NAT Gateway for private subnets.",
                    recommendation="Add a NAT Gateway to allow instances in private subnets to access the internet while remaining private."
                )
        
        # 9. Check for cost optimization
        if has_auto_scaling:
            add_best_practice(
                "cost",
                "Auto Scaling for Cost Optimization",
                "Architecture uses Auto Scaling to match capacity with demand, optimizing costs."
            )
        else:
            add_recommendation(
                "cost",
                "Implement Auto Scaling",
                "Consider implementing Auto Scaling to automatically adjust capacity based on demand, optimizing costs.",
                "high"
            )
        
        # 10. Check for disaster recovery components
        dr_components = [c for c in architecture_components if any(term in c.get("description", "").lower() for term in [
            "disaster recovery", "dr", "backup", "multi-region", "cross-region"
        ])]
        
        if not dr_components:
            add_recommendation(
                "reliability",
                "Disaster Recovery Planning",
                "Consider implementing disaster recovery components or multi-region deployment for critical workloads.",
                "medium"
            )
        
        # Generate overall validation summary
        total_issues = validation_results["summary"]["critical_issues"] + validation_results["summary"]["warning_issues"] + validation_results["summary"]["info_issues"]
        
        if validation_results["summary"]["critical_issues"] > 0:
            validation_results["validation_summary"] = f"Architecture has {validation_results['summary']['critical_issues']} critical issues that should be addressed immediately."
        elif validation_results["summary"]["warning_issues"] > 0:
            validation_results["validation_summary"] = f"Architecture has {validation_results['summary']['warning_issues']} warnings that should be reviewed."
        elif validation_results["summary"]["info_issues"] > 0:
            validation_results["validation_summary"] = f"Architecture has {validation_results['summary']['info_issues']} suggestions for improvement."
        else:
            validation_results["validation_summary"] = "Architecture validation completed successfully with no issues found."
        
        validation_results["validation_summary"] += f" {validation_results['summary']['best_practices_met']} best practices identified."
        
        return json.dumps(validation_results)
    except Exception as e:
        logger.error(f"Error validating architecture: {str(e)}")
        return json.dumps({
            "status": "error",
            "error": "Failed to validate architecture",
            "message": str(e)
        })

@tool
def check_aws_best_practices(architecture_components: List[Dict[str, Any]], category: Optional[str] = None) -> str:
    """Check AWS architecture against best practices in specific categories.

    Args:
        architecture_components (List[Dict[str, Any]]): List of architecture components with their relationships
        category (Optional[str]): Specific category to check (security, reliability, performance, cost, operational_excellence)
        
    Returns:
        str: JSON string containing best practice analysis results
    """
    if not architecture_components:
        return json.dumps({
            "status": "error",
            "error": "Empty architecture components",
            "message": "Please provide a list of architecture components."
        })
    
    try:
        # Initialize results
        results = {
            "status": "success",
            "summary": {
                "component_count": len(architecture_components),
                "categories_checked": [],
                "best_practices_met": 0,
                "best_practices_missing": 0
            },
            "best_practices_met": [],
            "best_practices_missing": [],
            "recommendations": []
        }
        
        # Determine which categories to check
        categories_to_check = []
        if category:
            if category.lower() in AWS_BEST_PRACTICES:
                categories_to_check = [category.lower()]
            else:
                return json.dumps({
                    "status": "error",
                    "error": "Invalid category",
                    "message": f"Category '{category}' not recognized. Valid categories are: security, reliability, performance, cost, operational_excellence"
                })
        else:
            categories_to_check = list(AWS_BEST_PRACTICES.keys())
        
        results["summary"]["categories_checked"] = categories_to_check
        
        # Helper function to add a met best practice
        def add_best_practice_met(category: str, practice: str, evidence: str):
            results["best_practices_met"].append({
                "category": category,
                "practice": practice,
                "evidence": evidence
            })
            results["summary"]["best_practices_met"] += 1
        
        # Helper function to add a missing best practice
        def add_best_practice_missing(category: str, practice: str, recommendation: str):
            results["best_practices_missing"].append({
                "category": category,
                "practice": practice,
                "recommendation": recommendation
            })
            results["summary"]["best_practices_missing"] += 1
        
        # Track components by type for easier reference
        components_by_type = {}
        for component in architecture_components:
            component_type = component.get("type", "").lower()
            if component_type not in components_by_type:
                components_by_type[component_type] = []
            components_by_type[component_type].append(component)
        
        # Check security best practices
        if "security" in categories_to_check:
            # Check for VPC usage
            if "vpc" in components_by_type:
                add_best_practice_met(
                    "security",
                    "Use security groups and NACLs to control traffic",
                    f"Architecture includes {len(components_by_type['vpc'])} VPC(s) for network isolation."
                )
            else:
                add_best_practice_missing(
                    "security",
                    "Use security groups and NACLs to control traffic",
                    "Add a VPC with proper security groups and NACLs to control network traffic."
                )
            
            # Check for IAM usage
            if any(t in components_by_type for t in ["iam", "cognito", "directory service"]):
                add_best_practice_met(
                    "security",
                    "Implement identity and access management",
                    "Architecture includes identity management services."
                )
            else:
                add_best_practice_missing(
                    "security",
                    "Implement identity and access management",
                    "Add IAM or Cognito for proper identity and access management."
                )
            
            # Check for encryption services
            if any(t in components_by_type for t in ["kms", "cloudhsm", "certificate manager"]):
                add_best_practice_met(
                    "security",
                    "Enable encryption for data at rest and in transit",
                    "Architecture includes encryption services."
                )
            else:
                add_best_practice_missing(
                    "security",
                    "Enable encryption for data at rest and in transit",
                    "Add KMS for encryption of data at rest and Certificate Manager for data in transit."
                )
            
            # Check for WAF
            if "waf" in components_by_type:
                add_best_practice_met(
                    "security",
                    "Use AWS WAF for web application protection",
                    "Architecture includes AWS WAF for web application protection."
                )
            elif any(t in components_by_type for t in ["alb", "cloudfront", "api gateway", "appsync"]):
                add_best_practice_missing(
                    "security",
                    "Use AWS WAF for web application protection",
                    "Add AWS WAF to protect your web applications from common exploits."
                )
            
            # Check for DDoS protection
            if "shield" in components_by_type:
                add_best_practice_met(
                    "security",
                    "Use AWS Shield for DDoS protection",
                    "Architecture includes AWS Shield for DDoS protection."
                )
            elif any(t in components_by_type for t in ["alb", "cloudfront", "api gateway"]):
                add_best_practice_missing(
                    "security",
                    "Use AWS Shield for DDoS protection",
                    "Add AWS Shield for DDoS protection of your public-facing resources."
                )
        
        # Check reliability best practices
        if "reliability" in categories_to_check:
            # Check for multi-AZ mentions
            has_multi_az = False
            for component in architecture_components:
                if "multi-az" in component.get("description", "").lower():
                    has_multi_az = True
                    break
            
            if has_multi_az:
                add_best_practice_met(
                    "reliability",
                    "Deploy across multiple Availability Zones",
                    "Architecture mentions multi-AZ deployment for high availability."
                )
            else:
                add_best_practice_missing(
                    "reliability",
                    "Deploy across multiple Availability Zones",
                    "Deploy critical components across multiple Availability Zones for high availability."
                )
            
            # Check for Auto Scaling
            if any(t in components_by_type for t in ["auto scaling", "autoscaling"]):
                add_best_practice_met(
                    "reliability",
                    "Use Auto Scaling for handling load changes",
                    "Architecture includes Auto Scaling for handling variable loads."
                )
            elif "ec2" in components_by_type:
                add_best_practice_missing(
                    "reliability",
                    "Use Auto Scaling for handling load changes",
                    "Add Auto Scaling to automatically adjust capacity based on demand."
                )
            
            # Check for health checks/monitoring
            if any(t in components_by_type for t in ["cloudwatch", "elb", "alb", "nlb", "route 53"]):
                add_best_practice_met(
                    "reliability",
                    "Implement health checks and monitoring",
                    "Architecture includes components that support health checks and monitoring."
                )
            else:
                add_best_practice_missing(
                    "reliability",
                    "Implement health checks and monitoring",
                    "Add health checks and monitoring to detect and recover from failures."
                )
            
            # Check for backup services
            if "backup" in components_by_type:
                add_best_practice_met(
                    "reliability",
                    "Implement backup and restore strategies",
                    "Architecture includes AWS Backup for centralized backup management."
                )
            elif any(t in components_by_type for t in ["rds", "dynamodb", "ebs", "efs", "s3"]):
                add_best_practice_missing(
                    "reliability",
                    "Implement backup and restore strategies",
                    "Add AWS Backup to centrally manage and automate backups of your data."
                )
        
        # Check performance best practices
        if "performance" in categories_to_check:
            # Check for CloudFront
            if "cloudfront" in components_by_type:
                add_best_practice_met(
                    "performance",
                    "Use CloudFront for content delivery",
                    "Architecture includes CloudFront for content delivery."
                )
            elif "s3" in components_by_type:
                add_best_practice_missing(
                    "performance",
                    "Use CloudFront for content delivery",
                    "Add CloudFront to improve content delivery performance and reduce latency."
                )
            
            # Check for caching
            if any(t in components_by_type for t in ["elasticache", "redis", "memcached", "cloudfront"]):
                add_best_practice_met(
                    "performance",
                    "Implement caching strategies",
                    "Architecture includes caching components."
                )
            else:
                add_best_practice_missing(
                    "performance",
                    "Implement caching strategies",
                    "Consider adding ElastiCache or CloudFront to implement caching and improve performance."
                )
            
            # Check for read replicas
            has_read_replicas = False
            for component in architecture_components:
                if "read replica" in component.get("description", "").lower():
                    has_read_replicas = True
                    break
            
            if has_read_replicas:
                add_best_practice_met(
                    "performance",
                    "Use read replicas for database read scaling",
                    "Architecture mentions read replicas for database read scaling."
                )
            elif any(t in components_by_type for t in ["rds", "aurora"]):
                add_best_practice_missing(
                    "performance",
                    "Use read replicas for database read scaling",
                    "Add read replicas to your databases to scale read operations and improve performance."
                )
        
        # Check cost optimization best practices
        if "cost" in categories_to_check:
            # Check for Auto Scaling (cost perspective)
            if any(t in components_by_type for t in ["auto scaling", "autoscaling"]):
                add_best_practice_met(
                    "cost",
                    "Use Auto Scaling to match capacity with demand",
                    "Architecture includes Auto Scaling to optimize resource usage and cost."
                )
            elif "ec2" in components_by_type:
                add_best_practice_missing(
                    "cost",
                    "Use Auto Scaling to match capacity with demand",
                    "Add Auto Scaling to automatically adjust capacity based on demand and optimize costs."
                )
            
            # Check for S3 lifecycle policies
            has_lifecycle_policies = False
            for component in architecture_components:
                if "lifecycle" in component.get("description", "").lower() and component.get("type", "").lower() == "s3":
                    has_lifecycle_policies = True
                    break
            
            if has_lifecycle_policies:
                add_best_practice_met(
                    "cost",
                    "Implement lifecycle policies for S3 objects",
                    "Architecture mentions S3 lifecycle policies for cost-effective storage management."
                )
            elif "s3" in components_by_type:
                add_best_practice_missing(
                    "cost",
                    "Implement lifecycle policies for S3 objects",
                    "Add lifecycle policies to your S3 buckets to automatically transition objects to lower-cost storage classes."
                )
            
            # Check for reserved instances mention
            has_reserved_instances = False
            for component in architecture_components:
                if "reserved" in component.get("description", "").lower():
                    has_reserved_instances = True
                    break
            
            if has_reserved_instances:
                add_best_practice_met(
                    "cost",
                    "Use Reserved Instances for predictable workloads",
                    "Architecture mentions Reserved Instances for cost savings."
                )
            elif "ec2" in components_by_type or "rds" in components_by_type:
                add_best_practice_missing(
                    "cost",
                    "Use Reserved Instances for predictable workloads",
                    "Consider using Reserved Instances for predictable workloads to reduce costs."
                )
        
        # Check operational excellence best practices
        if "operational_excellence" in categories_to_check:
            # Check for CloudFormation/IaC
            if "cloudformation" in components_by_type:
                add_best_practice_met(
                    "operational_excellence",
                    "Use Infrastructure as Code",
                    "Architecture includes CloudFormation for Infrastructure as Code."
                )
            else:
                add_best_practice_missing(
                    "operational_excellence",
                    "Use Infrastructure as Code",
                    "Consider using CloudFormation or AWS CDK to manage infrastructure as code."
                )
            
            # Check for monitoring and alerting
            if "cloudwatch" in components_by_type:
                add_best_practice_met(
                    "operational_excellence",
                    "Implement comprehensive monitoring and alerting",
                    "Architecture includes CloudWatch for monitoring and alerting."
                )
            else:
                add_best_practice_missing(
                    "operational_excellence",
                    "Implement comprehensive monitoring and alerting",
                    "Add CloudWatch for monitoring, alerting, and observability of your architecture."
                )
            
            # Check for Systems Manager
            if "systems manager" in components_by_type:
                add_best_practice_met(
                    "operational_excellence",
                    "Use AWS Systems Manager for operational management",
                    "Architecture includes Systems Manager for operational management."
                )
            elif "ec2" in components_by_type:
                add_best_practice_missing(
                    "operational_excellence",
                    "Use AWS Systems Manager for operational management",
                    "Add Systems Manager to automate operational tasks and improve visibility."
                )
            
            # Check for CloudTrail
            if "cloudtrail" in components_by_type:
                add_best_practice_met(
                    "operational_excellence",
                    "Use AWS CloudTrail for API activity monitoring",
                    "Architecture includes CloudTrail for API activity monitoring and auditing."
                )
            else:
                add_best_practice_missing(
                    "operational_excellence",
                    "Use AWS CloudTrail for API activity monitoring",
                    "Add CloudTrail to track API activity and changes to your AWS resources."
                )
        
        # Generate overall summary
        if results["summary"]["best_practices_met"] > 0 and results["summary"]["best_practices_missing"] == 0:
            results["overall_assessment"] = "Excellent! The architecture follows all checked AWS best practices."
        elif results["summary"]["best_practices_met"] > results["summary"]["best_practices_missing"]:
            results["overall_assessment"] = "Good. The architecture follows many AWS best practices, with some opportunities for improvement."
        elif results["summary"]["best_practices_met"] == results["summary"]["best_practices_missing"]:
            results["overall_assessment"] = "Fair. The architecture follows some AWS best practices, but has an equal number of opportunities for improvement."
        else:
            results["overall_assessment"] = "Needs improvement. The architecture has several opportunities to better align with AWS best practices."
        
        return json.dumps(results)
    except Exception as e:
        logger.error(f"Error checking AWS best practices: {str(e)}")
        return json.dumps({
            "status": "error",
            "error": "Failed to check AWS best practices",
            "message": str(e)
        })

@tool
def analyze_vpc_configuration(architecture_components: List[Dict[str, Any]]) -> str:
    """Analyze VPC configuration and validate correct placement of services inside/outside VPC.

    Args:
        architecture_components (List[Dict[str, Any]]): List of architecture components with their relationships
        
    Returns:
        str: JSON string containing VPC configuration analysis
    """
    if not architecture_components:
        return json.dumps({
            "status": "error",
            "error": "Empty architecture components",
            "message": "Please provide a list of architecture components."
        })
    
    try:
        # Initialize analysis results
        analysis_results = {
            "status": "success",
            "vpc_count": 0,
            "vpc_details": [],
            "service_placement": {
                "correctly_placed": [],
                "incorrectly_placed": []
            },
            "recommendations": []
        }
        
        # Find all VPCs
        vpcs = [c for c in architecture_components if c.get("type", "").lower() == "vpc"]
        analysis_results["vpc_count"] = len(vpcs)
        
        # If no VPCs found, provide recommendations
        if not vpcs:
            analysis_results["recommendations"].append({
                "priority": "high",
                "title": "Add VPC for network isolation",
                "description": "Your architecture doesn't include any VPC. Consider adding a VPC to isolate your resources and control network traffic."
            })
            
            # Check if there are services that should be in a VPC
            vpc_required_services = [c for c in architecture_components if c.get("type", "").lower() in VPC_INTERNAL_SERVICES]
            
            if vpc_required_services:
                service_names = [c.get("name", c.get("id", "Unknown")) for c in vpc_required_services]
                analysis_results["recommendations"].append({
                    "priority": "high",
                    "title": "Place services in VPC",
                    "description": f"The following services typically should be placed within a VPC: {', '.join(service_names)}"
                })
            
            return json.dumps(analysis_results)
        
        # Analyze each VPC
        for vpc in vpcs:
            vpc_id = vpc.get("id")
            vpc_name = vpc.get("name", vpc_id)
            
            # Find components in this VPC
            vpc_components = [c for c in architecture_components if c.get("vpc") == vpc_id]
            
            # Check for key networking components
            has_internet_gateway = any(c.get("type", "").lower() in ["internet gateway", "igw"] for c in vpc_components)
            has_nat_gateway = any(c.get("type", "").lower() in ["nat gateway", "nat"] for c in vpc_components)
            has_route_tables = any(c.get("type", "").lower() == "route table" for c in vpc_components)
            has_subnets = any(c.get("type", "").lower() == "subnet" for c in vpc_components)
            has_security_groups = any(c.get("type", "").lower() == "security group" for c in vpc_components)
            has_nacls = any(c.get("type", "").lower() in ["nacl", "network acl"] for c in vpc_components)
            
            # Analyze subnet configuration if present
            subnet_analysis = []
            if has_subnets:
                subnets = [c for c in vpc_components if c.get("type", "").lower() == "subnet"]
                public_subnets = []
                private_subnets = []
                
                for subnet in subnets:
                    subnet_name = subnet.get("name", subnet.get("id", "Unknown"))
                    subnet_description = subnet.get("description", "").lower()
                    
                    if "public" in subnet_name.lower() or "public" in subnet_description:
                        public_subnets.append(subnet_name)
                    elif "private" in subnet_name.lower() or "private" in subnet_description:
                        private_subnets.append(subnet_name)
                    else:
                        # If not explicitly labeled, try to infer from connections
                        subnet_id = subnet.get("id")
                        connected_to_igw = False
                        
                        for component in architecture_components:
                            if component.get("type", "").lower() in ["internet gateway", "igw"]:
                                if subnet_id in component.get("connections", []) or component.get("id") in subnet.get("connections", []):
                                    connected_to_igw = True
                                    break
                        
                        if connected_to_igw:
                            public_subnets.append(subnet_name)
                        else:
                            private_subnets.append(subnet_name)
                
                subnet_analysis = {
                    "public_subnets": public_subnets,
                    "private_subnets": private_subnets
                }
            
            # Add VPC details to results
            vpc_details = {
                "vpc_id": vpc_id,
                "vpc_name": vpc_name,
                "component_count": len(vpc_components),
                "networking_components": {
                    "has_internet_gateway": has_internet_gateway,
                    "has_nat_gateway": has_nat_gateway,
                    "has_route_tables": has_route_tables,
                    "has_subnets": has_subnets,
                    "has_security_groups": has_security_groups,
                    "has_nacls": has_nacls
                }
            }
            
            if subnet_analysis:
                vpc_details["subnet_analysis"] = subnet_analysis
            
            # Add recommendations for this VPC
            vpc_recommendations = []
            
            if not has_internet_gateway and (not subnet_analysis or subnet_analysis.get("public_subnets")):
                vpc_recommendations.append({
                    "priority": "medium",
                    "title": "Add Internet Gateway",
                    "description": f"VPC {vpc_name} has public subnets but no Internet Gateway. Add an Internet Gateway to allow communication with the internet."
                })
            
            if not has_nat_gateway and subnet_analysis and subnet_analysis.get("private_subnets") and subnet_analysis.get("public_subnets"):
                vpc_recommendations.append({
                    "priority": "medium",
                    "title": "Add NAT Gateway",
                    "description": f"VPC {vpc_name} has both public and private subnets but no NAT Gateway. Add a NAT Gateway to allow instances in private subnets to access the internet."
                })
            
            if not has_security_groups:
                vpc_recommendations.append({
                    "priority": "high",
                    "title": "Add Security Groups",
                    "description": f"VPC {vpc_name} does not have explicitly defined Security Groups. Add Security Groups to control inbound and outbound traffic at the instance level."
                })
            
            if vpc_recommendations:
                vpc_details["recommendations"] = vpc_recommendations
                analysis_results["recommendations"].extend(vpc_recommendations)
            
            analysis_results["vpc_details"].append(vpc_details)
        
        # Check for correct service placement
        for component in architecture_components:
            component_id = component.get("id")
            component_name = component.get("name", component_id)
            component_type = component.get("type", "").lower()
            vpc_id = component.get("vpc")
            
            # Skip VPCs themselves
            if component_type == "vpc":
                continue
            
            # Check if service should be in VPC but isn't
            if component_type in VPC_INTERNAL_SERVICES and not vpc_id:
                analysis_results["service_placement"]["incorrectly_placed"].append({
                    "component_id": component_id,
                    "component_name": component_name,
                    "component_type": component_type,
                    "issue": "Should be inside a VPC",
                    "recommendation": f"Move {component_name} inside a VPC for better security and network isolation."
                })
            # Check if service should be outside VPC but is inside
            elif component_type in VPC_EXTERNAL_SERVICES and vpc_id:
                analysis_results["service_placement"]["incorrectly_placed"].append({
                    "component_id": component_id,
                    "component_name": component_name,
                    "component_type": component_type,
                    "issue": "Should be outside a VPC",
                    "recommendation": f"{component_name} is a global or regional service and doesn't need to be placed inside a VPC."
                })
            else:
                analysis_results["service_placement"]["correctly_placed"].append({
                    "component_id": component_id,
                    "component_name": component_name,
                    "component_type": component_type,
                    "vpc_id": vpc_id
                })
        
        # Generate summary
        incorrect_count = len(analysis_results["service_placement"]["incorrectly_placed"])
        if incorrect_count > 0:
            analysis_results["summary"] = f"Found {incorrect_count} services with incorrect VPC placement. See details in service_placement.incorrectly_placed."
        else:
            analysis_results["summary"] = "All services are correctly placed with respect to VPC boundaries."
        
        return json.dumps(analysis_results)
    except Exception as e:
        logger.error(f"Error analyzing VPC configuration: {str(e)}")
        return json.dumps({
            "status": "error",
            "error": "Failed to analyze VPC configuration",
            "message": str(e)
        })