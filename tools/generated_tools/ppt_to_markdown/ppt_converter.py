"""
PPT to Markdown converter tool module.

This module provides tools for parsing PowerPoint (PPT/PPTX) files and converting
their content to Markdown format while preserving the hierarchical structure.
"""

import os
import json
import re
from typing import Dict, List, Any, Optional, Tuple, Union
from pathlib import Path
import tempfile
import logging
from datetime import datetime

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.shapes.base import BaseShape
    from pptx.shapes.autoshape import Shape
    from pptx.shapes.group import GroupShape
    from pptx.shapes.picture import Picture
    from pptx.shapes.graphfrm import GraphicFrame
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
except ImportError:
    raise ImportError(
        "Python-pptx library is required. Install it using: pip install python-pptx"
    )

from strands import tool

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s - %(name)s - %(levelname)s - %(message)s",
)
logger = logging.getLogger(__name__)


@tool
def ppt_to_markdown_converter(
    file_path: str,
    include_page_numbers: bool = True,
    include_slide_titles: bool = True,
    extract_notes: bool = False,
    include_images: bool = False,
    output_file: Optional[str] = None,
) -> str:
    """
    Convert a PowerPoint (PPT/PPTX) file to Markdown format.

    This tool parses a PowerPoint presentation and converts its content to Markdown format,
    preserving the hierarchical structure of the original presentation. It extracts text elements
    such as titles, paragraphs, and lists, and converts them to appropriate Markdown syntax.

    Args:
        file_path (str): Path to the PowerPoint file (.ppt or .pptx)
        include_page_numbers (bool, optional): Whether to include slide numbers in the output. Defaults to True.
        include_slide_titles (bool, optional): Whether to include slide titles as headers. Defaults to True.
        extract_notes (bool, optional): Whether to extract and include speaker notes. Defaults to False.
        include_images (bool, optional): Whether to include image references. Defaults to False.
        output_file (str, optional): Path to save the Markdown output. If None, only returns the result.

    Returns:
        str: JSON string containing the conversion result with the following structure:
            {
                "status": "success" or "error",
                "markdown_content": "The full Markdown content" (if successful),
                "error_message": "Error description" (if error occurred),
                "metadata": {
                    "slide_count": Number of slides processed,
                    "title": Presentation title (if available),
                    "processing_time": Time taken to process the file (in seconds)
                }
            }

    Raises:
        FileNotFoundError: If the specified file does not exist
        ValueError: If the file is not a valid PowerPoint file
        Exception: For other errors during processing
    """
    start_time = datetime.now()
    
    try:
        # Validate file existence
        if not os.path.exists(file_path):
            return json.dumps({
                "status": "error",
                "error_message": f"File not found: {file_path}",
                "metadata": {}
            })
        
        # Validate file extension
        file_ext = os.path.splitext(file_path)[1].lower()
        if file_ext not in ['.ppt', '.pptx']:
            return json.dumps({
                "status": "error",
                "error_message": f"Invalid file format: {file_ext}. Only .ppt and .pptx files are supported.",
                "metadata": {}
            })
        
        # Load the presentation
        logger.info(f"Loading presentation from {file_path}")
        presentation = Presentation(file_path)
        
        # Process the presentation
        markdown_content = []
        presentation_title = "Untitled Presentation"
        
        # Try to extract presentation title from the first slide if it has a title placeholder
        if presentation.slides and len(presentation.slides) > 0:
            first_slide = presentation.slides[0]
            for shape in first_slide.shapes:
                if shape.is_placeholder and shape.placeholder_format.type == 1:  # Title placeholder
                    if shape.has_text_frame and shape.text_frame.text:
                        presentation_title = shape.text_frame.text
                        break
        
        # Add presentation title as main header
        markdown_content.append(f"# {presentation_title}\n")
        
        # Process each slide
        for slide_index, slide in enumerate(presentation.slides):
            slide_number = slide_index + 1
            
            # Add slide separator with slide number
            if include_page_numbers:
                markdown_content.append(f"\n## Slide {slide_number}\n")
            
            # Extract slide title if available and requested
            slide_title = None
            if include_slide_titles:
                for shape in slide.shapes:
                    if shape.is_placeholder and shape.placeholder_format.type == 1:  # Title placeholder
                        if shape.has_text_frame and shape.text_frame.text:
                            slide_title = shape.text_frame.text
                            if include_page_numbers:
                                markdown_content.append(f"\n### {slide_title}\n")
                            else:
                                markdown_content.append(f"\n## {slide_title}\n")
                            break
            
            # Process all shapes in the slide
            for shape in slide.shapes:
                shape_markdown = _process_shape(shape, include_images)
                if shape_markdown:
                    markdown_content.append(shape_markdown)
            
            # Extract speaker notes if requested
            if extract_notes and slide.has_notes_slide and slide.notes_slide.notes_text_frame.text:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    markdown_content.append("\n> **Speaker Notes:**\n")
                    for line in notes_text.split('\n'):
                        markdown_content.append(f"> {line}\n")
        
        # Join all markdown content
        full_markdown = "\n".join(markdown_content)
        
        # Calculate processing time
        processing_time = (datetime.now() - start_time).total_seconds()
        
        # Save to file if output_file is specified
        if output_file:
            try:
                with open(output_file, 'w', encoding='utf-8') as f:
                    f.write(full_markdown)
                logger.info(f"Markdown content saved to {output_file}")
            except Exception as e:
                logger.error(f"Failed to save output file: {str(e)}")
                return json.dumps({
                    "status": "error",
                    "error_message": f"Failed to save output file: {str(e)}",
                    "markdown_content": full_markdown,
                    "metadata": {
                        "slide_count": len(presentation.slides),
                        "title": presentation_title,
                        "processing_time": processing_time
                    }
                })
        
        # Return the result
        return json.dumps({
            "status": "success",
            "markdown_content": full_markdown,
            "metadata": {
                "slide_count": len(presentation.slides),
                "title": presentation_title,
                "processing_time": processing_time
            }
        })
    
    except Exception as e:
        logger.error(f"Error converting PPT to Markdown: {str(e)}")
        return json.dumps({
            "status": "error",
            "error_message": f"Error converting PPT to Markdown: {str(e)}",
            "metadata": {
                "processing_time": (datetime.now() - start_time).total_seconds()
            }
        })


def _process_shape(shape: BaseShape, include_images: bool = False) -> str:
    """
    Process a PowerPoint shape and convert it to Markdown.
    
    Args:
        shape: The PowerPoint shape to process
        include_images: Whether to include image references
        
    Returns:
        str: Markdown representation of the shape
    """
    markdown = ""
    
    # Process text frame if available
    if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
        markdown += _process_text_frame(shape.text_frame)
    
    # Process table if available
    elif hasattr(shape, 'has_table') and shape.has_table:
        markdown += _process_table(shape.table)
    
    # Process chart if available
    elif hasattr(shape, 'has_chart') and shape.has_chart:
        chart_title = "Chart"
        if hasattr(shape.chart, 'chart_title') and shape.chart.chart_title:
            if shape.chart.chart_title.has_text_frame:
                chart_title = shape.chart.chart_title.text_frame.text or "Chart"
        markdown += f"\n*[Chart: {chart_title}]*\n"
    
    # Process image if available and requested
    elif include_images and (isinstance(shape, Picture) or 
                            (hasattr(shape, 'shape_type') and 
                             shape.shape_type == MSO_SHAPE_TYPE.PICTURE)):
        image_name = getattr(shape, 'name', 'Image')
        markdown += f"\n![{image_name}](image_{id(shape)})\n"
    
    # Process group shape
    elif isinstance(shape, GroupShape):
        for child_shape in shape.shapes:
            markdown += _process_shape(child_shape, include_images)
    
    return markdown


def _process_text_frame(text_frame) -> str:
    """
    Process a text frame and convert it to Markdown.
    
    Args:
        text_frame: The PowerPoint text frame
        
    Returns:
        str: Markdown representation of the text frame
    """
    if not text_frame.paragraphs:
        return ""
    
    markdown = []
    list_level = None
    in_list = False
    
    for i, paragraph in enumerate(text_frame.paragraphs):
        if not paragraph.text.strip():
            if not in_list:
                markdown.append("")  # Empty line for paragraph breaks
            continue
        
        # Check if this is a list item
        current_level = paragraph.level
        is_list_item = hasattr(paragraph, 'bullet') and paragraph.bullet.visible
        
        # Handle list formatting
        if is_list_item:
            in_list = True
            if list_level is None or current_level != list_level:
                list_level = current_level
            
            # Determine indentation based on level
            indent = "  " * current_level
            markdown.append(f"{indent}* {_format_paragraph_text(paragraph)}")
        else:
            in_list = False
            list_level = None
            
            # Regular paragraph formatting
            formatted_text = _format_paragraph_text(paragraph)
            
            # Check if it might be a header based on formatting
            if i == 0 and len(text_frame.paragraphs) > 1:
                # First paragraph might be a header if it's bold or larger font
                if any(run.font.bold for run in paragraph.runs if hasattr(run.font, 'bold')):
                    markdown.append(f"\n#### {formatted_text}\n")
                else:
                    markdown.append(formatted_text)
            else:
                markdown.append(formatted_text)
    
    return "\n".join(markdown)


def _format_paragraph_text(paragraph) -> str:
    """
    Format paragraph text with Markdown styling.
    
    Args:
        paragraph: The PowerPoint paragraph
        
    Returns:
        str: Formatted Markdown text
    """
    if not paragraph.runs:
        return paragraph.text.strip()
    
    formatted_parts = []
    
    for run in paragraph.runs:
        text = run.text.strip()
        if not text:
            continue
        
        # Apply formatting
        if hasattr(run.font, 'bold') and run.font.bold:
            text = f"**{text}**"
        if hasattr(run.font, 'italic') and run.font.italic:
            text = f"*{text}*"
        if hasattr(run.font, 'underline') and run.font.underline:
            text = f"<u>{text}</u>"
        
        formatted_parts.append(text)
    
    return " ".join(formatted_parts)


def _process_table(table) -> str:
    """
    Process a PowerPoint table and convert it to Markdown.
    
    Args:
        table: The PowerPoint table
        
    Returns:
        str: Markdown representation of the table
    """
    if not table.rows or not table.columns:
        return ""
    
    markdown = ["\n"]  # Start with a newline
    
    # Process each row
    for i, row in enumerate(table.rows):
        row_cells = []
        
        # Process each cell in the row
        for cell in row.cells:
            # Get cell text, replacing newlines with spaces
            cell_text = " ".join(cell.text_frame.text.split("\n"))
            row_cells.append(cell_text)
        
        # Add the row to markdown
        markdown.append("| " + " | ".join(row_cells) + " |")
        
        # Add header separator after first row
        if i == 0:
            markdown.append("| " + " | ".join(["---"] * len(row_cells)) + " |")
    
    markdown.append("")  # End with a newline
    return "\n".join(markdown)


@tool
def extract_ppt_structure(file_path: str) -> str:
    """
    Extract the structure of a PowerPoint presentation without converting to Markdown.
    
    This tool analyzes a PowerPoint file and returns a structured representation of its
    content, including slides, shapes, and text elements. This is useful for understanding
    the presentation structure before conversion.
    
    Args:
        file_path (str): Path to the PowerPoint file (.ppt or .pptx)
        
    Returns:
        str: JSON string containing the presentation structure with the following format:
            {
                "status": "success" or "error",
                "presentation": {
                    "slide_count": Number of slides,
                    "title": Presentation title (if available),
                    "slides": [
                        {
                            "slide_number": Slide number,
                            "title": Slide title (if available),
                            "content_elements": [
                                {
                                    "type": Element type (text, table, chart, image, etc.),
                                    "content": Element content or description
                                }
                            ]
                        }
                    ]
                },
                "error_message": Error description (if error occurred)
            }
    
    Raises:
        FileNotFoundError: If the specified file does not exist
        ValueError: If the file is not a valid PowerPoint file
        Exception: For other errors during processing
    """
    try:
        # Validate file existence
        if not os.path.exists(file_path):
            return json.dumps({
                "status": "error",
                "error_message": f"File not found: {file_path}"
            })
        
        # Validate file extension
        file_ext = os.path.splitext(file_path)[1].lower()
        if file_ext not in ['.ppt', '.pptx']:
            return json.dumps({
                "status": "error",
                "error_message": f"Invalid file format: {file_ext}. Only .ppt and .pptx files are supported."
            })
        
        # Load the presentation
        presentation = Presentation(file_path)
        
        # Extract presentation structure
        presentation_title = "Untitled Presentation"
        
        # Try to extract presentation title from the first slide if it has a title placeholder
        if presentation.slides and len(presentation.slides) > 0:
            first_slide = presentation.slides[0]
            for shape in first_slide.shapes:
                if shape.is_placeholder and shape.placeholder_format.type == 1:  # Title placeholder
                    if shape.has_text_frame and shape.text_frame.text:
                        presentation_title = shape.text_frame.text
                        break
        
        # Process slides
        slides_data = []
        for slide_index, slide in enumerate(presentation.slides):
            slide_number = slide_index + 1
            slide_data = {
                "slide_number": slide_number,
                "title": None,
                "content_elements": []
            }
            
            # Extract slide title if available
            for shape in slide.shapes:
                if shape.is_placeholder and shape.placeholder_format.type == 1:  # Title placeholder
                    if shape.has_text_frame and shape.text_frame.text:
                        slide_data["title"] = shape.text_frame.text
                        break
            
            # Process all shapes in the slide
            for shape in slide.shapes:
                element = _extract_shape_info(shape)
                if element:
                    slide_data["content_elements"].append(element)
            
            # Add speaker notes if available
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    slide_data["content_elements"].append({
                        "type": "speaker_notes",
                        "content": notes_text
                    })
            
            slides_data.append(slide_data)
        
        # Return the result
        return json.dumps({
            "status": "success",
            "presentation": {
                "slide_count": len(presentation.slides),
                "title": presentation_title,
                "slides": slides_data
            }
        })
    
    except Exception as e:
        logger.error(f"Error extracting PPT structure: {str(e)}")
        return json.dumps({
            "status": "error",
            "error_message": f"Error extracting PPT structure: {str(e)}"
        })


def _extract_shape_info(shape: BaseShape) -> Dict[str, Any]:
    """
    Extract information about a PowerPoint shape.
    
    Args:
        shape: The PowerPoint shape
        
    Returns:
        Dict: Information about the shape
    """
    # Text frame
    if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
        text_content = shape.text_frame.text.strip()
        if text_content:
            return {
                "type": "text",
                "content": text_content
            }
    
    # Table
    elif hasattr(shape, 'has_table') and shape.has_table:
        table_data = []
        for row in shape.table.rows:
            row_data = []
            for cell in row.cells:
                row_data.append(cell.text_frame.text)
            table_data.append(row_data)
        
        return {
            "type": "table",
            "content": table_data
        }
    
    # Chart
    elif hasattr(shape, 'has_chart') and shape.has_chart:
        chart_title = "Chart"
        if hasattr(shape.chart, 'chart_title') and shape.chart.chart_title:
            if shape.chart.chart_title.has_text_frame:
                chart_title = shape.chart.chart_title.text_frame.text or "Chart"
        
        return {
            "type": "chart",
            "content": chart_title
        }
    
    # Picture
    elif isinstance(shape, Picture) or (hasattr(shape, 'shape_type') and 
                                       shape.shape_type == MSO_SHAPE_TYPE.PICTURE):
        image_name = getattr(shape, 'name', 'Image')
        return {
            "type": "image",
            "content": image_name
        }
    
    # Group shape
    elif isinstance(shape, GroupShape):
        group_elements = []
        for child_shape in shape.shapes:
            child_info = _extract_shape_info(child_shape)
            if child_info:
                group_elements.append(child_info)
        
        if group_elements:
            return {
                "type": "group",
                "content": group_elements
            }
    
    # Other shapes
    elif hasattr(shape, 'shape_type'):
        return {
            "type": "shape",
            "content": f"Shape (Type: {shape.shape_type})"
        }
    
    return None


@tool
def batch_convert_ppt_to_markdown(
    directory_path: str,
    output_directory: Optional[str] = None,
    include_page_numbers: bool = True,
    include_slide_titles: bool = True,
    extract_notes: bool = False,
    include_images: bool = False,
) -> str:
    """
    Convert multiple PowerPoint files in a directory to Markdown format.
    
    This tool processes all PowerPoint files (.ppt and .pptx) in the specified directory
    and converts them to Markdown format. It creates one Markdown file for each PowerPoint file.
    
    Args:
        directory_path (str): Path to the directory containing PowerPoint files
        output_directory (str, optional): Directory to save Markdown files. If None, uses the input directory.
        include_page_numbers (bool, optional): Whether to include slide numbers in the output. Defaults to True.
        include_slide_titles (bool, optional): Whether to include slide titles as headers. Defaults to True.
        extract_notes (bool, optional): Whether to extract and include speaker notes. Defaults to False.
        include_images (bool, optional): Whether to include image references. Defaults to False.
        
    Returns:
        str: JSON string containing the conversion results with the following structure:
            {
                "status": "success" or "error",
                "conversion_results": [
                    {
                        "input_file": Path to the input PPT file,
                        "output_file": Path to the output Markdown file,
                        "status": "success" or "error",
                        "error_message": Error description (if error occurred)
                    }
                ],
                "summary": {
                    "total_files": Total number of PowerPoint files found,
                    "successful_conversions": Number of successful conversions,
                    "failed_conversions": Number of failed conversions
                },
                "error_message": Error description (if error occurred)
            }
    
    Raises:
        FileNotFoundError: If the specified directory does not exist
        Exception: For other errors during processing
    """
    try:
        # Validate directory existence
        if not os.path.exists(directory_path) or not os.path.isdir(directory_path):
            return json.dumps({
                "status": "error",
                "error_message": f"Directory not found or not a directory: {directory_path}",
                "conversion_results": [],
                "summary": {
                    "total_files": 0,
                    "successful_conversions": 0,
                    "failed_conversions": 0
                }
            })
        
        # Set output directory
        if output_directory is None:
            output_directory = directory_path
        else:
            # Create output directory if it doesn't exist
            os.makedirs(output_directory, exist_ok=True)
        
        # Find all PowerPoint files in the directory
        ppt_files = []
        for file in os.listdir(directory_path):
            file_path = os.path.join(directory_path, file)
            if os.path.isfile(file_path) and file.lower().endswith(('.ppt', '.pptx')):
                ppt_files.append(file_path)
        
        # Process each file
        conversion_results = []
        successful_conversions = 0
        failed_conversions = 0
        
        for file_path in ppt_files:
            file_name = os.path.basename(file_path)
            output_file_name = os.path.splitext(file_name)[0] + '.md'
            output_file_path = os.path.join(output_directory, output_file_name)
            
            try:
                # Convert the file
                result_json = ppt_to_markdown_converter(
                    file_path=file_path,
                    include_page_numbers=include_page_numbers,
                    include_slide_titles=include_slide_titles,
                    extract_notes=extract_notes,
                    include_images=include_images,
                    output_file=output_file_path
                )
                
                result = json.loads(result_json)
                
                if result["status"] == "success":
                    successful_conversions += 1
                    conversion_results.append({
                        "input_file": file_path,
                        "output_file": output_file_path,
                        "status": "success"
                    })
                else:
                    failed_conversions += 1
                    conversion_results.append({
                        "input_file": file_path,
                        "output_file": output_file_path,
                        "status": "error",
                        "error_message": result.get("error_message", "Unknown error")
                    })
            
            except Exception as e:
                failed_conversions += 1
                conversion_results.append({
                    "input_file": file_path,
                    "output_file": output_file_path,
                    "status": "error",
                    "error_message": str(e)
                })
        
        # Return the summary
        return json.dumps({
            "status": "success",
            "conversion_results": conversion_results,
            "summary": {
                "total_files": len(ppt_files),
                "successful_conversions": successful_conversions,
                "failed_conversions": failed_conversions
            }
        })
    
    except Exception as e:
        logger.error(f"Error in batch conversion: {str(e)}")
        return json.dumps({
            "status": "error",
            "error_message": f"Error in batch conversion: {str(e)}",
            "conversion_results": [],
            "summary": {
                "total_files": 0,
                "successful_conversions": 0,
                "failed_conversions": 0
            }
        })